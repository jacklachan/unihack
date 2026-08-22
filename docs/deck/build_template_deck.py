"""Fill the mandatory UniHack prototype template with CALIPER's content.

The template is used exactly as supplied: every slide keeps its background
artwork, its heading text box, and its position. This script only writes text
into the empty body area of each slide and drops in screenshots of the running
prototype. Nothing in the template's design is moved, restyled, or removed.

    python docs/deck/build_template_deck.py

Reads  docs/deck/template/unihack_template_user.pptx
Writes docs/deck/CALIPER_UniHack.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "template", "unihack_template_user.pptx")
SHOTS = os.path.join(HERE, "shots")
OUT = os.path.join(HERE, "CALIPER_UniHack.pptx")

# Drawn from the template's own header bar, so added content belongs to it.
NAVY = RGBColor(0x00, 0x38, 0x6F)
BLUE = RGBColor(0x0B, 0x5F, 0xA5)
INK = RGBColor(0x1B, 0x24, 0x32)
MUTE = RGBColor(0x5A, 0x66, 0x78)
BRASS = RGBColor(0x8A, 0x64, 0x17)
RULE = RGBColor(0xD5, 0xDE, 0xE8)
WASH = RGBColor(0xF3, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x30, 0x21)

BODY = "Calibri"
HEAD = "Calibri"

# The white content area between the template's header and footer bars.
L, T, W = 0.52, 1.58, 8.96
BOTTOM = 5.16


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return box, tf


def para(tf, text, size=12, bold=False, color=INK, first=False, space=6,
         font=BODY, align=PP_ALIGN.LEFT, italic=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def rich(tf, parts, size=12, first=False, space=6, line=1.18):
    """One paragraph built from (text, bold, colour) runs."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    p.line_spacing = line
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = BODY
    return p


def card(slide, l, t, w, h, fill=WASH, line=RULE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def stat(slide, l, t, w, value, label, color=BLUE):
    _, tf = tb(slide, l, t, w, 0.72)
    para(tf, value, size=25, bold=True, color=color, first=True, space=0)
    para(tf, label, size=9.5, color=MUTE, space=0)


def node(slide, l, t, w, h, title, sub="", fill=WASH, edge=RULE,
         tcolor=NAVY, scolor=MUTE):
    card(slide, l, t, w, h, fill, edge)
    _, tf = tb(slide, l + 0.10, t + 0.10, w - 0.20, h - 0.20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, title, size=10.5, bold=True, color=tcolor, first=True, space=1,
         align=PP_ALIGN.CENTER, line=1.0)
    if sub:
        para(tf, sub, size=8.5, color=scolor, space=0, align=PP_ALIGN.CENTER,
             line=1.0)


def arrow(slide, l, t, w, h=0.16, color=BLUE, shape=MSO_SHAPE.RIGHT_ARROW):
    sh = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w),
                                Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def picture(slide, path, l, t, max_w, max_h, border=True):
    """Place a screenshot scaled to fit the given box, centred horizontally."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    left = l + (max_w - w) / 2.0
    pic = slide.shapes.add_picture(path, Inches(left), Inches(t),
                                   Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = RULE
        pic.line.width = Pt(0.75)
    return pic


def heading_text(slide):
    """The template's own heading box on this slide, if it has one."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh
    return None


def caption(slide, t, text):
    _, tf = tb(slide, L, t, W, 0.3)
    para(tf, text, size=9, color=MUTE, first=True, space=0, italic=True)


# --------------------------------------------------------------------------

prs = Presentation(TEMPLATE)
S = prs.slides

# ---- Slide 2 : Team details ---------------------------------------------
sl = S[1]
box = heading_text(sl)
tf = box.text_frame
# Keep the template's own labels; make the blanks impossible to miss.
for p in tf.paragraphs:
    txt = p.text.strip()
    if txt.startswith("Team name"):
        for r in p.runs:
            r.text = ""
        p.runs[0].text = "Team name:  "
        r = p.add_run()
        r.text = "<<< FILL IN >>>"
        r.font.color.rgb = RED
        r.font.bold = True
        r.font.size = p.runs[0].font.size
    elif txt.startswith("Team leader"):
        for r in p.runs:
            r.text = ""
        p.runs[0].text = "Team leader name:  "
        r = p.add_run()
        r.text = "<<< FILL IN >>>"
        r.font.color.rgb = RED
        r.font.bold = True
        r.font.size = p.runs[0].font.size

# ---- Slide 3 : Brief about your solution --------------------------------
sl = S[2]
_, tf = tb(sl, L, T, W, 1.15)
rich(tf, [("CALIPER turns a messy supplier row into Unilog's ", False, INK),
          ("252-column delivery format", True, INK),
          (" — and can show you why it wrote every single cell.", False, INK)],
     size=13.5, first=True, space=7)
rich(tf, [("Most pipelines hand the row to a language model and validate what "
           "comes back. CALIPER makes that structurally impossible. Every "
           "extractor — rules, the brand registry, the taxonomy, family "
           "consensus and the model — writes into one typed ", False, INK),
          ("Product Fact Graph, where a fact without evidence is rejected at "
           "insertion", True, NAVY),
          (". Composition, validation and export only ever read from it. The "
           "model is never permitted to write an output cell: it may propose a "
           "fact that quotes the source, or judge a fact that already exists.",
           False, INK)], size=11.5, space=0)

for i, (v, l) in enumerate([("1,000 → 252", "rows to columns, in ~3 seconds"),
                            ("100%", "INVOICE_DESC within 40 characters"),
                            ("99.5%", "sibling agreement, 1,516 comparisons"),
                            ("26,596", "cells carrying their own evidence")]):
    stat(sl, L + i * 2.24, 3.05, 2.2, v, l)

_, tf = tb(sl, L, 3.95, W, 1.0)
rich(tf, [("Runs with no API key and no third-party packages. ", True, NAVY),
          ("A judge opens the link and 1,000 real industrial SKUs are already "
           "enriched; clicking any row reveals the rule that produced each "
           "value and the exact characters of the input that justified it.",
           False, INK)], size=11.5, first=True, space=0)

# ---- Slide 4 : How it enriches minimal product information ---------------
sl = S[3]
box = heading_text(sl)
tf = box.text_frame
keep = tf.paragraphs[0].text
for p in list(tf.paragraphs)[1:]:
    p._element.getparent().remove(p._element)

_, tf = tb(sl, L, T - 0.22, W, 3.5)
rich(tf, [("Input is six columns", True, NAVY),
          (" — part number, description, three inconsistent brand fields and a "
           "manufacturer string. From that, CALIPER derives facts and never "
           "invents them:", False, INK)], size=11.5, first=True, space=8)

for head, text in [
    ("Parse, don't guess.",
     "Dimensions, grits, wattages, counts and materials are lifted from the "
     "description by typed rules. Each becomes a fact carrying the exact "
     "substring it came from: a description reading 6 in x .045 in x 7/8 in "
     "yields three separate measurements, each traceable."),
    ("Resolve identity against a registry.",
     "\"Milw\" becomes Milwaukee®; distributor noise like \"(4031)\" is "
     "stripped. Where brand and manufacturer disagree, that is reported as a "
     "defect in the supplier's data rather than silently overwritten."),
    ("Borrow from siblings, with corroboration.",
     "Products are clustered into families by part-number and description "
     "shape. A fact missing on one row is filled only when at least two "
     "agreeing siblings supply it — and disagreement lowers confidence."),
    ("Induce the category rulebook from the data.",
     "Which attributes a category has, and which are filterable, is derived "
     "from the rows themselves with a fill rate behind every attribute — the "
     "work normally done by hand for tens of thousands of categories."),
    ("Compose to the format, not to a prompt.",
     "Five descriptions are generated from the same facts under their own "
     "length and casing rules, so they cannot contradict each other."),
]:
    rich(tf, [(head + "  ", True, BLUE), (text, False, INK)], size=10.5,
         space=5)

# ---- Slide 5 : Opportunities --------------------------------------------
sl = S[4]
_, tf = tb(sl, L, T + 0.85, W, 3.0)
for q, a in [
    ("How different is it from existing ideas?",
     "Every other approach we found prompts a model and validates afterwards, "
     "so a wrong value has already been written. CALIPER inverts the control "
     "flow: the model cannot reach an output cell at all. It proposes facts "
     "that must quote the source, or audits facts that already exist — and an "
     "unquoted proposal is discarded, not corrected."),
    ("How does it solve the problem statement?",
     "It produces all 252 delivery columns from six, at ~3 seconds per 1,000 "
     "rows, with 77.4% of rows ready to publish unattended and the remainder "
     "routed to a review queue ranked by how much that review is worth. It "
     "also reports defects it finds in the supplier's own catalogue."),
    ("USP of the proposed solution",
     "Auditability as an architectural guarantee, not a report. 26,596 "
     "populated cells each carry a rule id, a confidence and the source "
     "characters behind them. Because two labelled rows cannot support a "
     "claim, we also publish a label-free measure — sibling agreement across "
     "all 1,000 rows, 99.5% over 1,516 comparisons."),
]:
    para(tf, q, size=12, bold=True, color=NAVY, first=(q.startswith("How diff")),
         space=3)
    para(tf, a, size=10.5, color=INK, space=11, line=1.16)

# ---- Slide 6 : Features --------------------------------------------------
sl = S[5]
feats = [
    ("Evidence panel", "Click a row: every populated column shows its method, "
                       "rule id, confidence and the source substring."),
    ("Character-budget solver", "INVOICE_DESC's 40-character upper-case "
                                "ceiling solved as a budget problem. 100% "
                                "compliant by construction."),
    ("Schema detection", "Column roles are detected, not assumed — the same "
                         "pipeline runs on a file whose headers it has never "
                         "seen."),
    ("Category-spec induction", "Unilog-style category rulebooks derived from "
                                "the rows, with a fill rate behind every "
                                "attribute."),
    ("Family consensus + anomaly flags", "Gaps filled only from corroborating "
                                         "siblings; rows that break their "
                                         "family pattern are surfaced."),
    ("Physical guardrails", "Domain ranges and dimensional coherence. These "
                            "caught a bug in our own parser."),
    ("Findings on the source data", "Brand/manufacturer conflicts and "
                                    "unresolvable rows reported rather than "
                                    "absorbed."),
    ("Review queue", "The rows a human should look at, ranked by the value of "
                     "looking."),
    ("Optional model, firewalled", "Bring your own key for extra recall. Held "
                                   "in memory only; proposals without a "
                                   "verbatim quote are discarded."),
]
cw, ch, gx, gy = 2.88, 0.90, 0.16, 0.14
for i, (t_, d) in enumerate(feats):
    r, c = divmod(i, 3)
    l = L + c * (cw + gx)
    t = T + 0.30 + r * (ch + gy)
    card(sl, l, t, cw, ch)
    _, tf = tb(sl, l + 0.12, t + 0.11, cw - 0.24, ch - 0.22)
    para(tf, t_, size=10, bold=True, color=NAVY, first=True, space=2, line=1.0)
    para(tf, d, size=8.5, color=INK, space=0, line=1.1)

# ---- Slide 7 : Process flow ---------------------------------------------
sl = S[6]
steps = [
    ("1 · Ingest", "CSV or XLSX read\nwith the standard library"),
    ("2 · Detect schema", "Column roles inferred\nfrom headers and values"),
    ("3 · Extract facts", "Typed rules, registry, taxonomy.\nNo evidence, "
                          "no fact"),
    ("4 · Corroborate", "Family clustering; gaps filled\nonly by agreeing "
                        "siblings"),
    ("5 · Compose", "Five descriptions from one\nfact set, under their limits"),
    ("6 · Validate & route", "Guardrails, then ready\nor review queue"),
]
bw, bh = 2.72, 0.86
xs = [L + 0.06, L + 0.06 + bw + 0.42, L + 0.06 + 2 * (bw + 0.42)]
for i, (t_, s_) in enumerate(steps):
    r, c = divmod(i, 3)
    t = T + 0.26 + r * 1.42
    node(sl, xs[c], t, bw, bh, t_, s_)
    if c < 2:
        arrow(sl, xs[c] + bw + 0.09, t + bh / 2 - 0.08, 0.24)
# wrap from step 3 down to step 4
arrow(sl, xs[2] + bw / 2 - 0.08, T + 0.26 + bh + 0.14, 0.16, 0.30,
      shape=MSO_SHAPE.DOWN_ARROW)

_, tf = tb(sl, L, T + 3.02, W, 0.75)
rich(tf, [("The one-way rule: ", True, NAVY),
          ("steps 1–4 write facts into the graph. Steps 5–6 only read from it. "
           "A model may join step 3 as a proposer whose quote is checked "
           "against the source, or step 6 as an auditor — it is never given "
           "write access to an output column.", False, INK)],
     size=10.5, first=True, space=0)

# ---- Slide 8 : Wireframes / prototype layout ----------------------------
sl = S[7]
picture(sl, os.path.join(SHOTS, "panel_top.png"), L, T + 0.10, 4.42, 3.16)
_, tf = tb(sl, L + 4.66, T + 0.14, W - 4.66, 3.1)
para(tf, "The interface, in the order a judge meets it", size=12, bold=True,
     color=NAVY, first=True, space=8)
for k, v in [
    ("It is already running.",
     "The page enriches 1,000 SKUs on load. No button, no key, no setup — the "
     "first ten seconds show output, not a form."),
    ("Upload, or use the sample.",
     "Drop in any catalogue CSV. Column roles are detected and printed back."),
    ("The model is opt-in.",
     "Ticking the box reveals a provider and a key field. Left alone, the "
     "deterministic path still fills all 252 columns."),
    ("Then: catalogue, evidence, findings, induced specs, downloads.",
     "Five tabs, in that order — the evidence panel sits directly under the "
     "table so one click explains a row."),
]:
    rich(tf, [(k + "  ", True, BLUE), (v, False, INK)], size=9.5, space=7)
caption(sl, 4.90, "Screenshot of the running prototype, light theme; the page "
                  "also renders correctly in dark mode.")

# ---- Slide 9 : Architecture ---------------------------------------------
sl = S[8]
top = T + 0.26
# Writers
node(sl, L, top, 2.30, 0.34, "Rule extractors", fill=WHITE)
node(sl, L, top + 0.42, 2.30, 0.34, "Brand / manufacturer registry", fill=WHITE)
node(sl, L, top + 0.84, 2.30, 0.34, "Taxonomy + induced specs", fill=WHITE)
node(sl, L, top + 1.26, 2.30, 0.34, "Family consensus", fill=WHITE)
node(sl, L, top + 1.68, 2.30, 0.34, "Model proposer (quotes only)",
     fill=RGBColor(0xFD, 0xF4, 0xE2), edge=RGBColor(0xE0, 0xC9, 0x94),
     tcolor=BRASS)
_, tf = tb(sl, L, top - 0.26, 2.30, 0.24)
para(tf, "WRITE", size=8.5, bold=True, color=MUTE, first=True, space=0,
     align=PP_ALIGN.CENTER)

for i in range(5):
    arrow(sl, L + 2.36, top + 0.09 + i * 0.42, 0.30, 0.16)

# The graph
card(sl, L + 2.72, top, 3.10, 2.02, RGBColor(0xEC, 0xF3, 0xFA), NAVY)
_, tf = tb(sl, L + 2.84, top + 0.14, 2.86, 1.76)
para(tf, "PRODUCT FACT GRAPH", size=11.5, bold=True, color=NAVY, first=True,
     space=4, align=PP_ALIGN.CENTER)
para(tf, "Typed facts, each with evidence,\nconfidence and provenance",
     size=9, color=INK, space=7, align=PP_ALIGN.CENTER, line=1.1)
para(tf, "A fact without evidence is\nrejected at insertion", size=9.5,
     bold=True, color=RED, space=7, align=PP_ALIGN.CENTER, line=1.1)
para(tf, "Agreement raises confidence.\nConflict lowers it.", size=9,
     color=MUTE, space=0, align=PP_ALIGN.CENTER, line=1.1)

for i in range(4):
    arrow(sl, L + 5.88, top + 0.30 + i * 0.42, 0.30, 0.16)

# Readers
node(sl, L + 6.24, top + 0.24, 2.34, 0.34, "Description composer", fill=WHITE)
node(sl, L + 6.24, top + 0.66, 2.34, 0.34, "Validator + guardrails", fill=WHITE)
node(sl, L + 6.24, top + 1.08, 2.34, 0.34, "252-column exporter", fill=WHITE)
node(sl, L + 6.24, top + 1.50, 2.34, 0.34, "Evidence panel / review queue",
     fill=WHITE)
_, tf = tb(sl, L + 6.24, top - 0.26, 2.34, 0.24)
para(tf, "READ ONLY", size=8.5, bold=True, color=MUTE, first=True, space=0,
     align=PP_ALIGN.CENTER)

_, tf = tb(sl, L, top + 2.22, W, 0.85)
rich(tf, [("The boundary is one-way and that is the whole design. ", True,
           NAVY),
          ("Readers cannot create values; writers cannot skip evidence. A "
           "model sits on the left as a proposer whose quote is verified "
           "against the source, and on the right as an auditor returning "
           "supported / unsupported / unknown — in neither role can it "
           "originate a cell.", False, INK)], size=10.5, first=True, space=0)

# ---- Slide 10 : Technologies --------------------------------------------
sl = S[9]
_, tf = tb(sl, L, T + 0.28, W, 0.5)
rich(tf, [("Python 3.12, standard library only. ", True, NAVY),
          ("The pipeline imports nothing outside the stdlib — CSV and XLSX are "
           "read and written with zipfile and XML directly. There is no "
           "pandas, no openpyxl, no ML framework and no dependency that can "
           "break a build.", False, INK)], size=11, first=True, space=0)

tech = [
    ("Core pipeline", "Python 3.12 · standard library only · no packages"),
    ("Data I/O", "csv, zipfile, xml.etree — XLSX written by hand"),
    ("Prototype UI", "Gradio 6 on Hugging Face Spaces (the one dependency)"),
    ("Optional model layer", "Provider-neutral adapter: Groq, Gemini, "
                             "Anthropic, OpenAI"),
    ("Model governance", "Evidence firewall, redundancy guard, verdict-only "
                         "auditing"),
    ("Quality assurance", "46 invariant tests; evaluation against the "
                          "published answer key"),
]
for i, (k, v) in enumerate(tech):
    r, c = divmod(i, 2)
    l = L + c * 4.56
    t = T + 0.98 + r * 0.72
    card(sl, l, t, 4.40, 0.62)
    _, tf = tb(sl, l + 0.14, t + 0.10, 4.12, 0.44)
    para(tf, k, size=10.5, bold=True, color=NAVY, first=True, space=1,
         line=1.0)
    para(tf, v, size=9, color=INK, space=0, line=1.05)

_, tf = tb(sl, L, T + 3.30, W, 0.4)
para(tf, "Consequence: a judge can clone the repository and run the whole "
         "pipeline with one command, on a machine with no network access.",
     size=10, color=MUTE, first=True, space=0, italic=True)

# ---- Slide 11 : Estimated implementation cost ---------------------------
sl = S[10]
_, tf = tb(sl, L, T + 0.22, W, 0.44)
rich(tf, [("Running the deterministic pipeline costs nothing but CPU.", True,
           NAVY)], size=11.5, first=True, space=0)

rowsdef = [
    ("Deterministic enrichment", "₹0", "No API calls. ~3 s per 1,000 rows on "
                                       "one core; ~1 CPU-hour per 1M SKUs."),
    ("Hosting the prototype", "₹0", "Hugging Face Spaces free tier."),
    ("Optional model pass", "≈ ₹40–90 per 10,000 rows",
     "Only rows the rules cannot resolve are sent — roughly 12% of the "
     "catalogue — and results are cached on disk."),
    ("Human review", "Scales with the queue", "77.4% of rows need no human. "
     "The queue is ranked, so reviewer time goes to the rows where it pays."),
]
t = T + 0.78
for k, v, d in rowsdef:
    card(sl, L, t, W, 0.72)
    _, tf = tb(sl, L + 0.16, t + 0.09, 3.0, 0.54)
    para(tf, k, size=10.5, bold=True, color=NAVY, first=True, space=1, line=1.0)
    _, tf = tb(sl, L + 3.20, t + 0.09, 1.9, 0.54)
    para(tf, v, size=11, bold=True, color=BLUE, first=True, space=0, line=1.0)
    _, tf = tb(sl, L + 5.16, t + 0.09, W - 5.32, 0.56)
    para(tf, d, size=9, color=INK, first=True, space=0, line=1.1)
    t += 0.80

# ---- Slide 12 : Snapshots of the MVP ------------------------------------
sl = S[11]
picture(sl, os.path.join(SHOTS, "evi_closeup.png"), L, T + 0.08, W, 2.56)
_, tf = tb(sl, L, T + 2.76, W, 1.05)
rich(tf, [("Three columns of one row, as the prototype explains them. ", True,
           NAVY),
          ("Each shows the method (rule, registry), the rule id, a confidence, "
           "and in the gold quote the exact characters of the input that "
           "justified the value — with the column they came from. "
           "BRAND_NAME became Milwaukee® because the four characters "
           "“milw” appear in Part_Desc and resolve against the approved "
           "brand registry; nothing here was generated.", False, INK)],
     size=10.5, first=True, space=6)
rich(tf, [("26,596 populated cells across the catalogue carry provenance of "
           "this kind, produced by the pipeline rather than written afterwards "
           "as a report. Reproduce this panel by clicking any row at "
           "huggingface.co/spaces/jacklachan/unihack.", False, MUTE)],
     size=9.5, space=0)
# (No separate caption: it collided with the paragraph above.)

# ---- Slide 13 : Future development --------------------------------------
sl = S[12]
_, tf = tb(sl, L, T + 0.20, W, 3.2)
rich(tf, [("What the measurements told us to build next.", True, NAVY)],
     size=11.5, first=True, space=9)
for head, text in [
    ("Manufacturer retrieval is the real ceiling.",
     "Eleven of the twelve attribute values in the labelled row appear nowhere "
     "in its input. No amount of prompting recovers them — they have to be "
     "fetched from manufacturer sources and then held to the same evidence "
     "rule as everything else. This is the single highest-value extension."),
    ("A larger answer key.",
     "Accuracy is 48.0% exact on two labelled rows. That base is too narrow to "
     "carry a claim, which is why we report sibling agreement alongside it. "
     "With a few hundred labelled rows the same harness reports a real number."),
    ("Learning from corrections.",
     "Every reviewer edit is already captured with the fact it overrode. "
     "Feeding those back as registry and vocabulary entries makes the "
     "deterministic path absorb the reviewer's judgement permanently."),
    ("Image and document evidence.",
     "The Fact/Evidence type does not care whether a quote came from text or "
     "from a spec sheet. Extending the evidence kind lets datasheets and "
     "images become first-class sources without weakening the rule."),
]:
    rich(tf, [(head + "  ", True, BLUE), (text, False, INK)], size=10.5,
         space=8)

# ---- Slide 14 : Links ----------------------------------------------------
sl = S[13]
links = [
    ("Working prototype", "https://huggingface.co/spaces/jacklachan/unihack",
     "Opens already enriched. No login, no key, no setup."),
    ("GitHub repository", "https://github.com/jacklachan/unihack",
     "Standard library only, 46 invariant tests, clone and run."),
    ("Demo video (3 minutes)", "<<< PASTE YOUR VIDEO LINK >>>",
     "Set sharing to anyone-with-the-link and check it in a private window."),
]
t = T + 1.00
for k, v, d in links:
    card(sl, L, t, W, 0.76)
    _, tf = tb(sl, L + 0.18, t + 0.09, W - 0.36, 0.60)
    para(tf, k, size=11, bold=True, color=NAVY, first=True, space=3, line=1.0)
    para(tf, v, size=11.5, bold=True,
         color=RED if v.startswith("<<<") else BLUE, space=3, line=1.0)
    para(tf, d, size=9, color=MUTE, space=0, line=1.0)
    t += 0.86

prs.save(OUT)
print("wrote", OUT)
