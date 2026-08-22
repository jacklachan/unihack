"""Build the CALIPER prototype deck.

    python docs/deck/build_deck.py

Palette is the project's own -- brass on deep indigo. Brass because a caliper
is a brass instrument and the whole thesis is measurement. Dark title and
close, light content between them.

Motif: brass-numbered markers and bordered stat blocks with monospaced
numerals. Deliberately no accent stripes and no rules under titles.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── tokens ──────────────────────────────────────────────────────────────
INK      = RGBColor(0x12, 0x18, 0x2B)
INK_LINE = RGBColor(0x2A, 0x34, 0x50)
PAPER    = RGBColor(0xF2, 0xF4, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BRASS    = RGBColor(0xB9, 0x8A, 0x34)
BRASS_D  = RGBColor(0x8A, 0x64, 0x17)
BRASS_BG = RGBColor(0xFB, 0xF4, 0xE6)
MUTE     = RGBColor(0x6B, 0x74, 0x88)
MUTE_D   = RGBColor(0x4A, 0x54, 0x68)
MUTE_L   = RGBColor(0x9A, 0xA3, 0xB2)
RULE     = RGBColor(0xD3, 0xD9, 0xE3)
GOOD     = RGBColor(0x1F, 0x7A, 0x4C)
FLAG     = RGBColor(0x9A, 0x43, 0x1F)
CODE_FG  = RGBColor(0xD8, 0xDE, 0xEA)
CODE_DIM = RGBColor(0x6E, 0x7A, 0x92)
CODE_BLU = RGBColor(0x8F, 0xA8, 0xC4)
DARK_SUB = RGBColor(0xC6, 0xCE, 0xDC)
DARK_MUT = RGBColor(0x85, 0x92, 0xA6)

SERIF, SANS, MONO = "Cambria", "Calibri", "Courier New"
W, H, M = 13.333, 7.5, 0.72

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── helpers ─────────────────────────────────────────────────────────────
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = INK if dark else PAPER
    return s


def text(s, txt, x, y, w, h, *, font=SANS, size=14, bold=False, italic=False,
         color=INK, spacing=None, align=PP_ALIGN.LEFT, char_space=None):
    """One text box. `txt` is a string or a list of (text, overrides) runs."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    chunks = txt if isinstance(txt, list) else [(txt, {})]

    # Runs share a paragraph until a newline appears. Starting a new paragraph
    # per chunk put "CALI" and "PER" on separate lines and dropped the second
    # one through the subtitle beneath it.
    para = tf.paragraphs[0]
    para.alignment = align
    if spacing:
        para.line_spacing = Pt(spacing)

    for chunk, over in chunks:
        parts = str(chunk).split("\n")
        for i, part in enumerate(parts):
            if i:
                para = tf.add_paragraph()
                para.alignment = over.get("align", align)
                if spacing:
                    para.line_spacing = Pt(spacing)
            if part == "":
                continue
            r = para.add_run()
            r.text = part
            f = r.font
            f.name = over.get("font", font)
            f.size = Pt(over.get("size", size))
            f.bold = over.get("bold", bold)
            f.italic = over.get("italic", italic)
            f.color.rgb = over.get("color", color)
            cs = over.get("char_space", char_space)
            if cs:
                r.font._rPr.set("spc", str(int(cs * 100)))
    return box


def rect(s, x, y, w, h, fill=WHITE, line=RULE, line_w=1):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.text_frame.text = ""
    return sh


def hairline(s, x, y, w, color=RULE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Emu(9525))          # 0.01"
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def marker(s, x, y, n):
    """Small brass disc with a number -- the repeated motif."""
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                            Inches(0.32), Inches(0.32))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = BRASS
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.name, r.font.size, r.font.bold = MONO, Pt(11), True
    r.font.color.rgb = INK
    return sh


def eyebrow(s, label, color=BRASS):
    text(s, label.upper(), M, 0.42, W - 2 * M, 0.3,
         font=MONO, size=11, bold=True, color=color, char_space=2.2)


def title(s, txt, color=INK, size=32):
    text(s, txt, M, 0.8, W - 2 * M, 1.0,
         font=SERIF, size=size, bold=True, color=color, spacing=size * 1.14)


def subtitle(s, txt, color=MUTE, y=1.72, h=0.62):
    text(s, txt, M, y, W - 2 * M, h, font=SANS, size=14.5, color=color,
         spacing=20)


def code(s, x, y, w, h, lines, size=11.5, spacing=17):
    rect(s, x, y, w, h, fill=INK, line=INK_LINE)
    text(s, lines, x + 0.24, y + 0.2, w - 0.48, h - 0.4,
         font=MONO, size=size, color=CODE_FG, spacing=spacing)


def stat(s, x, y, w, h, value, label, note=None, color=BRASS_D):
    rect(s, x, y, w, h)
    text(s, value, x + 0.22, y + 0.2, w - 0.44, 0.66,
         font=MONO, size=28, bold=True, color=color)
    text(s, label.upper(), x + 0.22, y + 0.92, w - 0.44, 0.3,
         font=SANS, size=10, bold=True, color=MUTE, char_space=1.0)
    if note:
        text(s, note, x + 0.22, y + 1.24, w - 0.44, 0.56,
             font=SANS, size=10.5, color=MUTE, spacing=13)


def footer(s, txt, color=MUTE_L):
    text(s, txt, M, H - 0.58, W - 2 * M, 0.32, font=SANS, size=10, color=color)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ═══════════════════════════════════════════════════════ 1 · TITLE
s = slide(dark=True)
text(s, "UniHack  ·  AI-Powered Product Intelligence for Industrial Commerce",
     M, 0.6, 11, 0.3, font=MONO, size=11, bold=True, color=BRASS, char_space=2)
text(s, [("CALI", {"color": WHITE}), ("PER", {"color": BRASS})],
     M, 1.5, 10, 1.5, font=SERIF, size=72, bold=True, char_space=4)
text(s, "An enrichment pipeline that measures what it claims,\n"
        "and refuses to write what it cannot support.",
     M, 3.2, 8.6, 1.2, font=SANS, size=18, color=DARK_SUB, spacing=27)
for i, (v, l) in enumerate([("1,000", "rows in 9s"), ("252", "columns out"),
                            ("26,627", "traceable cells"), ("42", "tests passing")]):
    x = M + i * 3.02
    text(s, v, x, 5.15, 2.8, 0.6, font=MONO, size=28, bold=True, color=BRASS)
    text(s, l.upper(), x, 5.8, 2.8, 0.3, font=SANS, size=10, bold=True,
         color=DARK_MUT, char_space=1.0)
notes(s, "CALIPER. The name is the thesis: a caliper is a measuring instrument, "
         "and every number in this deck was produced by a run rather than typed in.")

# ═══════════════════════════════════════════════════════ 2 · PROBLEM
s = slide()
eyebrow(s, "The problem")
title(s, "Six columns in. Two hundred and fifty-two out.")
subtitle(s, "One real row from the working dataset. Three of its five fields are "
            "placeholders that mean “empty”.")
code(s, M, 2.45, 6.3, 2.05, [
    ("Mfg_Part_Num : 49-94-1940\n", {}),
    ("Part_Desc    : 49-94-1940 Milw 14\"x1/8\"x1\"\n", {"color": WHITE}),
    ("               Masonry Cut Off Disc\n", {"color": WHITE}),
    ("E1_Brand     : -- Unbranded --\n", {"color": CODE_DIM}),
    ("DIB_Brand    : -- No DIB Brand --\n", {"color": CODE_DIM}),
    ("Part_Manuf   : Milwaukee Accessory (4031)", {}),
])
for i, (hd, sub) in enumerate([
        ("An approved brand", "exact casing, with the ® symbol"),
        ("A classpath", "the key every attribute is validated against"),
        ("Five descriptions", "five different length and casing rules"),
        ("Ordered attributes", "drawn from a controlled vocabulary")]):
    y = 2.45 + i * 0.64
    marker(s, 7.35, y + 0.02, i + 1)
    text(s, hd, 7.87, y - 0.02, 5.0, 0.3, size=14, bold=True)
    text(s, sub, 7.87, y + 0.26, 5.0, 0.3, size=11.5, color=MUTE)
rect(s, M, 4.95, W - 2 * M, 0.78, fill=RGBColor(0xF7, 0xEC, 0xE6), line=FLAG)
text(s, "The brief is blunt about the obvious approach: “a fluent description "
        "made of invented values scores zero.”",
     M + 0.28, 5.16, W - 2 * M - 0.56, 0.4, size=14, italic=True, color=FLAG)
footer(s, "Sample-1000_Items · one row, unedited")
notes(s, "The obvious approach is to hand the row to a language model and hope. "
         "The brief tells you exactly why that fails.")

# ═══════════════════════════════════════════════════════ 3 · THE IDEA
s = slide()
eyebrow(s, "The design")
title(s, "Facts before prose.")
subtitle(s, "Every extractor writes into one structure. Nothing downstream may "
            "invent a value — it can only render facts that already exist.")
code(s, M, 2.4, 7.4, 2.9, [
    ("brand         Milwaukee®             IDN-BRD-01  0.88\n", {"color": BRASS}),
    ("manufacturer  Milwaukee Tool         IDN-MFR-01  0.86\n", {}),
    ("dimensions    14 in x 1/8 in x 1 in  DIM-CHN-01  0.92\n", {}),
    ("diameter      14 in                  DIM-CHN-02  0.92\n", {}),
    ("thickness     1/8 in                 DIM-CHN-02  0.92\n", {}),
    ("arbor_size    1 in                   DIM-CHN-02  0.92\n", {}),
    ("application   Masonry                APP-MAT-01  0.92\n", {}),
    ("item_type     Cut Off Disc           ITM-LEX-01  0.95\n", {"color": BRASS}),
    ("\n", {}),
    ("→ Milwaukee® 49-94-1940 Cut Off Disc,\n", {"color": CODE_BLU}),
    ("  14 in x 1/8 in x 1 in, Masonry", {"color": CODE_BLU}),
], size=10.5, spacing=16)
for i, (hd, body) in enumerate([
        ("Nothing ungrounded gets in",
         "A fact without evidence is rejected at insertion. A tested invariant, not a convention."),
        ("The five descriptions agree",
         "They are five renderings of one fact set, not five separate generations."),
        ("Provenance is free",
         "26,627 cells, each traceable to the rule that made it and the characters behind it.")]):
    y = 2.42 + i * 1.0
    marker(s, 8.45, y + 0.02, i + 1)
    text(s, hd, 8.97, y - 0.02, 3.95, 0.3, size=13.5, bold=True)
    text(s, body, 8.97, y + 0.28, 3.95, 0.68, size=11, color=MUTE, spacing=13.5)
text(s, "Rules, registries, taxonomy, family consensus and the model all write "
        "facts.   Composition, validation and export only read them.",
     M, 5.55, W - 2 * M, 0.4, size=14, bold=True)
footer(s, "The one-way boundary is the entire design")
notes(s, "The model is never allowed to write an output cell. It proposes facts "
         "that must quote the source, or renders verdicts on facts that exist.")

# ═══════════════════════════════════════════════════════ 4 · FINDINGS
s = slide()
eyebrow(s, "What measurement changed")
title(s, "Four findings that came from measuring, not reasoning.")
subtitle(s, "Each one changed the build. Two contradicted decisions already made.")
FIND = [
    ("01", "The answer key fills 25% of the sheet",
     "UPC, UNSPSC, every dimension, SDS and country of origin are blank in both "
     "published rows. A pipeline filling 90% is inventing two-thirds of a catalogue."),
    ("02", "Attribute values are not in the supplier row",
     "The labelled dishwasher's input is 39 characters. Eleven of its twelve "
     "attribute values appear nowhere in it — they need the manufacturer's source."),
    ("03", "A guardrail caught our own parser",
     "Four wheels reported an arbor hole wider than the disc. The dimension regex "
     "knew inches but not millimetres. Findings went 4 → 2 → 0; both defects were real."),
    ("04", "A bug no reading could find",
     "A heredoc turned \\b into a literal backspace byte inside two regexes. grep, "
     "sed and cat all showed the files as correct — terminals don't draw it."),
]
for i, (num, hd, body) in enumerate(FIND):
    col, row = i % 2, i // 2
    x, y = M + col * 6.12, 2.52 + row * 1.78
    rect(s, x, y, 5.82, 1.6)
    text(s, num, x + 0.24, y + 0.2, 0.7, 0.34, font=MONO, size=15, bold=True,
         color=BRASS)
    text(s, hd, x + 0.92, y + 0.18, 4.66, 0.32, size=13, bold=True)
    text(s, body, x + 0.92, y + 0.54, 4.66, 0.94, size=11, color=MUTE,
         spacing=14)
footer(s, "A fifth killed a feature: family amortisation was estimated at 8× "
          "and measured at 1.77×")
notes(s, "Finding 2 is the important one. It sets an honest ceiling on attribute "
         "accuracy without retrieval.")

# ═══════════════════════════════════════════════════════ 5 · RESULTS
s = slide()
eyebrow(s, "Results")
title(s, "1,000 rows, nine seconds, no API key required.")
ROW1 = [("100%", "Invoice ≤ 40 chars", "compliant by construction, not by checking", GOOD),
        ("92.5%", "Brand resolved", "to an approved name with its ® symbol", BRASS_D),
        ("88.8%", "Classified", "the remainder abstain rather than guess", BRASS_D),
        ("77.4%", "Ready to publish", "774 rows need no human at all", GOOD)]
ROW2 = [("99.5%", "Sibling agreement", "over 1,516 comparisons, 185 families", BRASS_D),
        ("614", "Relationship edges", "powers · variant_of · same_series", BRASS_D),
        ("64", "Category specs", "induced from rows, not hand-written", BRASS_D),
        ("48.0%", "vs. the answer key", "on 2 labelled rows — a narrow base, stated", BRASS_D)]
for r, data in enumerate([ROW1, ROW2]):
    for i, (v, l, n, c) in enumerate(data):
        stat(s, M + i * 3.06, 2.3 + r * 2.0, 2.86, 1.86, v, l, n, c)
footer(s, "Every figure produced by a run and cross-checked against report.json "
          "before it was written down")
notes(s, "Invoice compliance is 100% because it is solved as a budget problem "
         "rather than asked of a model.")

# ═══════════════════════════════════════════════════════ 6 · CONSISTENCY
s = slide()
eyebrow(s, "Answering the obvious objection")
title(s, "Two labelled rows is a narrow base. Consistency is not.")
subtitle(s, "No care widens a calibration set of two — but self-consistency can be "
            "asked of all 1,000 rows: products in one family must agree about the "
            "facts they share.", h=0.7)
code(s, M, 2.8, 6.4, 2.1, [
    ("families with 2+ members  :   185\n", {}),
    ("attribute comparisons     : 1,516\n", {}),
    ("siblings agree            : ", {}),
    ("99.5%\n", {"color": BRASS, "bold": True}),
    ("\n", {}),
    ("brand 100%     manufacturer 100%\n", {"color": CODE_BLU}),
    ("item_type 99.5%   classpath 99.4%", {"color": CODE_BLU}),
], size=12, spacing=18)
rect(s, 7.55, 2.8, 5.06, 2.1, fill=BRASS_BG, line=BRASS)
text(s, "It earned its place within the hour.", 7.79, 3.0, 4.6, 0.3,
     size=13, bold=True)
text(s, "Brand agreement came back at 95.3%, which reads as an extraction error. "
        "It was a clustering error — every “Dishwasher SS” from one appliance "
        "co-op landed in a single family, because the brand lives in the part "
        "number, not the description.\n\nMaking family membership brand-aware "
        "took brand and manufacturer to 100%.",
     7.79, 3.36, 4.6, 1.44, size=10.5, color=MUTE_D, spacing=13)
text(s, "This is deliberately not called accuracy — a pipeline can be uniformly "
        "wrong and perfectly consistent. But an extractor that disagrees with "
        "itself across products differing only by size is certainly wrong somewhere.",
     M, 5.2, W - 2 * M, 0.6, size=13, spacing=18)
footer(s, "A quality signal measured on the whole catalogue rather than on two rows")
notes(s, "This is the counterweight to the small ground truth, and it immediately "
         "found a real defect.")

# ═══════════════════════════════════════════════════════ 7 · PROTOTYPE
s = slide()
eyebrow(s, "The prototype")
title(s, "Upload a catalogue. Click any cell. See why.")
subtitle(s, "A judge's own file works — column names are detected, not assumed. "
            "Proven on a test file with different headers, different order and two "
            "junk columns.")
PANES = [
    ("Evidence panel", "Every value names the rule that produced it, the characters "
     "of the input that justified it, and its confidence."),
    ("Budget solver trace", "The 40-character invoice line, showing which facts were "
     "kept, which abbreviated, which dropped."),
    ("Findings on the source", "Brand and manufacturer disagreeing; rows where "
     "classification refused to guess."),
    ("Induced specs", "64 category rulebooks derived from the rows, with a count "
     "behind every attribute."),
    ("Review queue", "Ranked by uncertainty × sibling SKUs affected — the best use "
     "of the next human minute."),
    ("Relationship graph", "614 edges: which battery powers which tool, which board "
     "varies from which."),
]
for i, (hd, body) in enumerate(PANES):
    col, row = i % 3, i // 3
    x, y = M + col * 4.07, 2.68 + row * 1.66
    rect(s, x, y, 3.82, 1.48)
    text(s, hd, x + 0.22, y + 0.2, 3.4, 0.3, size=13, bold=True, color=BRASS_D)
    text(s, body, x + 0.22, y + 0.54, 3.4, 0.82, size=10.5, color=MUTE, spacing=13)
text(s, "Exports the 252-column delivery file as CSV and XLSX, plus the provenance "
        "ledger, review queue and relationship graph.",
     M, 6.2, W - 2 * M, 0.35, size=12, italic=True, color=MUTE)
footer(s, "python -m caliper serve  ·  no install step, no API key, runs offline")
notes(s, "Everything here loads from the shipped files on a fresh clone. The "
         "evidence panel does not need a live run.")

# ═══════════════════════════════════════════════════════ 8 · DIFFERENT
s = slide()
eyebrow(s, "Why this one is different")
title(s, "Most pipelines generate, then check. This one cannot generate.")
COMPARE = [
    ("The usual approach", "CALIPER"),
    ("Model writes the cell; a validator checks it afterwards",
     "Model may only propose a fact that quotes the source — composition renders, it cannot invent"),
    ("A confidence score invented per rule",
     "Confidence rises from independent methods agreeing, and is checked against sibling consistency"),
    ("Fill as much of the sheet as possible",
     "Fill rate calibrated against the answer key, because over-filling is fabrication"),
    ("Character limits checked, sometimes missed",
     "A budget problem with abbreviation ladders — 100% compliant by construction"),
    ("A review queue that forgets its answers",
     "Corrections persist, scoped to part / family / category / brand, and replay every run"),
]
for i, (a, b) in enumerate(COMPARE):
    y = 2.34 + i * 0.72
    head = i == 0
    if not head:
        hairline(s, M, y - 0.1, W - 2 * M)
    text(s, a, M, y, 5.4, 0.62,
         size=10.5 if head else 12, bold=head, char_space=1.0 if head else None,
         color=MUTE if head else RGBColor(0x5A, 0x64, 0x78), spacing=15)
    text(s, b, 6.5, y, 6.1, 0.62,
         size=10.5 if head else 12, bold=True,
         char_space=1.0 if head else None,
         color=BRASS_D if head else INK, spacing=15)
footer(s, "No third-party packages · CSV and XLSX read and written with the "
          "standard library")
notes(s, "Hallucination is not defended against with a checker; it is made "
         "structurally impossible for the output cell.")

# ═══════════════════════════════════════════════════════ 9 · LIMITS
s = slide()
eyebrow(s, "What we would not claim", FLAG)
title(s, "The honest edges.")
subtitle(s, "A submission whose thesis is measurement has to report what it "
            "cannot measure.")
LIMITS = [
    ("Two labelled rows",
     "Every accuracy figure rests on a calibration set of two. The harness prints "
     "the sample size beside every rate and scales untouched to a larger file."),
    ("No manufacturer retrieval",
     "Finding 02 shows this, not a missing rule, is the binding constraint on "
     "attribute-value accuracy."),
    ("No official vocabularies",
     "The LOV, the 27,000-row manufacturer master and the UOM standards were never "
     "distributed. Our registry is a ~75-entry bootstrap, and LOV coverage is "
     "measured against an induced vocabulary and labelled as such."),
    ("Conformance is not quality",
     "100% character compliance is true, and is not the same as a good line. No "
     "automated check we have closes that gap."),
]
for i, (hd, body) in enumerate(LIMITS):
    y = 2.52 + i * 1.0
    marker(s, M, y + 0.02, i + 1)
    text(s, hd, M + 0.5, y - 0.02, 3.1, 0.3, size=13.5, bold=True)
    text(s, body, M + 3.7, y - 0.02, 8.5, 0.86, size=11.5, color=MUTE, spacing=14.5)
footer(s, "Roadmap: manufacturer-source retrieval and document intelligence — "
          "where the remaining columns live")
notes(s, "Naming the ceiling is stronger than hiding it, and finding 2 proves "
         "where it sits.")

# ═══════════════════════════════════════════════════════ 10 · CLOSE
s = slide(dark=True)
text(s, "Enrichment is table stakes.", M, 1.9, 11.6, 0.8,
     font=SERIF, size=34, bold=True, color=WHITE)
# Sized to break in exactly two lines. At 38pt this wrapped to three and the
# last one ran through the divider below it.
text(s, "Knowing which rows ship unattended —\n"
        "and being able to prove why — is the product.",
     M, 2.72, 11.9, 1.5, font=SERIF, size=34, bold=True, color=BRASS, spacing=44)
hairline(s, M, 4.72, W - 2 * M, INK_LINE)
for i, (l, v) in enumerate([("Prototype", "python -m caliper serve"),
                            ("Repository", "github.com/jacklachan/unihack"),
                            ("Written brief", "docs/submission.html")]):
    x = M + i * 4.07
    text(s, l.upper(), x, 4.95, 3.9, 0.28, size=10, bold=True,
         color=RGBColor(0x7C, 0x87, 0x98), char_space=1.2)
    text(s, v, x, 5.25, 3.9, 0.32, font=MONO, size=12, color=WHITE)
text(s, "Every number in this deck was produced by a run — then the outputs "
        "behind it were read.",
     M, 6.4, 11.2, 0.3, size=11.5, italic=True, color=RGBColor(0x7C, 0x87, 0x98))
notes(s, "Close on the thesis: the differentiator is not that we enrich, it is "
         "that we can say which rows are safe to publish unattended, and show "
         "the evidence for each one.")

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "CALIPER_UniHack.pptx")
prs.save(out)
print("wrote {} ({} slides)".format(out, len(prs.slides.__iter__.__self__._sldIdLst)))
